from pptx import Presentation
from os.path import expanduser
import os
import json


# Create a presentation object with arguments path, note name and
class Powerpoint(object):
    ''' slides to be converted to raw text '''

    def __init__(self, path, title, filetype):
        self.path = path
        self.title = title
        self.raw_text = []
        self.prs = Presentation(path)
        self.filetype = filetype

    # method that will extract text from the
    # slides of a presentation and store it in a list 
    def slide_note(self):

        # list containing all text and counter for
        # slide number
        slide_counter = 1

        # loop through each slide
        for slide in self.prs.slides:

            # keep count of the slide number when extracting text
            slide_number = str(slide_counter)

            # list containing our raw text for a slide
            text_runs = []
            if self.filetype == 'org':
                text_runs.append('#Slide '+slide_number+'\n'+'\n')
            if self.filetype == 'latex':
                text_runs.append('%Slide '+slide_number+'\n'+'\n')
            slide_counter += 1

            # loop through all detected shapes whithin a slide
            for shape in slide.shapes:

                # skip frame if no text is found
                if not shape.has_text_frame:
                    continue

                # retriwve text from each paragraph and add to textrun list
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)

            # add newlines after slides, join into a single string
            # add slide text string to rawtext list
            text_runs.append('\n'+'\n')
            extracted_text = ''.join(text_runs)
            self.raw_text.append(extracted_text)

    # method that will convert raw text into a .tex or .org file
    def note_file(self):
        if self.filetype == 'latex':
            ftp = ".tex"
        if self.filetype == 'org':
            ftp = ".org"


        path = expanduser('~/.local/lib/python2.7/site-packages/pynoter/markup.json')
        newfile = self.title + ftp
        hbt_data = path

        # load data head, body and tail strings
        with open(hbt_data) as f:
            markup = json.load(f)

        for dict in markup:
            if dict['language'] == self.filetype:
                head = dict['head']
                body = dict['body']
                tail = dict['tail']
            
        # create a new file with rawtext head body and tail
        with open(newfile, 'w') as fobj:
            fobj.write(head)
            for slide in self.raw_text:
                fobj.write(slide.encode('utf8'))
                fobj.write(body)
            fobj.write(tail)
            
        
        


